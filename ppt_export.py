"""
PowerPoint export helpers.

Build per-tab PowerPoint decks where each section is its own slide,
generously spaced, with Team Saudi branding.

Usage:
    sections = [
        {"title": "Medal Table",  "kind": "table", "df": medals_df},
        {"title": "By Sport",     "kind": "chart", "fig": fig},
        {"title": "Note",         "kind": "text",  "body": "..."},
    ]
    ppt_bytes = build_pptx("Overview", sections)
"""
from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt


ELITE      = RGBColor(0x23, 0x50, 0x36)
ENABLER    = RGBColor(0x69, 0xc3, 0x99)
DISCIPLINE = RGBColor(0x18, 0x34, 0x2a)
STAMINA    = RGBColor(0xc3, 0xd9, 0xd1)
VICTORY    = RGBColor(0xeb, 0xce, 0x83)
WHITE      = RGBColor(0xff, 0xff, 0xff)

SLIDE_W = Inches(13.333)   # 16:9 widescreen
SLIDE_H = Inches(7.5)


def _add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _add_header(slide, title: str, logo_path: Path | None = None):
    """Green header strip with title text and optional logo on the right."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0), SLIDE_W, Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = ELITE
    bar.line.fill.background()

    # Gold accent stripe below the header
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0.85), SLIDE_W, Inches(0.06))
    accent.fill.solid(); accent.fill.fore_color.rgb = VICTORY
    accent.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(11), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE

    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(11.5), Inches(0.1),
                                 height=Inches(0.65))

    # Footer green stripe
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(7.30), SLIDE_W, Inches(0.20))
    footer.fill.solid(); footer.fill.fore_color.rgb = ELITE
    footer.line.fill.background()


def _add_cover(prs, deck_title: str, subtitle: str, logo_path: Path | None):
    slide = _add_blank_slide(prs)

    # Full-bleed Elite Green background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = ELITE
    bg.line.fill.background()

    # Discipline-green band across the middle for the title block
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), SLIDE_W, Inches(2.9))
    band.fill.solid(); band.fill.fore_color.rgb = DISCIPLINE
    band.line.fill.background()

    # Victory-gold accent stripe at the top of the band
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.25), SLIDE_W, Inches(0.08))
    accent.fill.solid(); accent.fill.fore_color.rgb = VICTORY
    accent.line.fill.background()

    # Team Saudi logo (larger, prominent)
    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(0.7), Inches(0.6), height=Inches(1.4))

    # Title block centred in the band
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(12), Inches(2.5))
    p = tx.text_frame.paragraphs[0]
    p.text = "GCC GAMES DOHA 2026"
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE

    p2 = tx.text_frame.add_paragraph()
    p2.text = deck_title
    p2.font.size = Pt(28); p2.font.bold = True; p2.font.color.rgb = VICTORY
    p2.space_before = Pt(10)

    p3 = tx.text_frame.add_paragraph()
    p3.text = subtitle
    p3.font.size = Pt(14); p3.font.color.rgb = STAMINA
    p3.space_before = Pt(12)

    # Footer credit
    foot = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(12), Inches(0.4))
    pf = foot.text_frame.paragraphs[0]
    pf.text = "Team Saudi · Performance Analysis"
    pf.font.size = Pt(11); pf.font.color.rgb = STAMINA; pf.font.italic = True


def _add_table_slide(prs, title: str, df: pd.DataFrame, max_rows: int = 18,
                     logo_path: Path | None = None):
    slide = _add_blank_slide(prs)
    _add_header(slide, title, logo_path)

    if df is None or df.empty:
        tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        tx.text_frame.text = "(no data)"
        return

    df = df.head(max_rows).fillna("").astype(str)
    rows, cols = df.shape

    # Position the table with generous margins
    left, top = Inches(0.6), Inches(1.3)
    width, height = Inches(12.1), Inches(min(5.8, 0.4 + 0.35 * rows))
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Header row
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid(); cell.fill.fore_color.rgb = ELITE
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = WHITE; run.font.bold = True; run.font.size = Pt(11)

    # Data rows
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10); run.font.color.rgb = DISCIPLINE
            if i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = STAMINA

    # Footnote
    if rows < df.shape[0]:
        note = slide.shapes.add_textbox(Inches(0.6), Inches(7), Inches(12), Inches(0.4))
        p = note.text_frame.paragraphs[0]
        p.text = f"Showing top {max_rows} rows."
        p.font.size = Pt(9); p.font.color.rgb = DISCIPLINE


def _fig_to_png(fig, width: int = 1200, height: int = 600) -> bytes | None:
    """Render a Plotly figure to PNG bytes. Returns None if kaleido isn't installed."""
    try:
        return fig.to_image(format="png", width=width, height=height, scale=2)
    except Exception:
        return None


def _add_chart_slide(prs, title: str, fig, logo_path: Path | None = None):
    slide = _add_blank_slide(prs)
    _add_header(slide, title, logo_path)
    img = _fig_to_png(fig) if fig is not None else None
    if img:
        bio = BytesIO(img)
        slide.shapes.add_picture(bio, Inches(0.7), Inches(1.2),
                                 width=Inches(11.9))
    else:
        tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
        p = tx.text_frame.paragraphs[0]
        p.text = "(chart export needs the `kaleido` package — run `pip install kaleido` and retry)"
        p.font.size = Pt(14); p.font.color.rgb = DISCIPLINE


def _add_text_slide(prs, title: str, body: str, logo_path: Path | None = None):
    slide = _add_blank_slide(prs)
    _add_header(slide, title, logo_path)
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
    for line in (body or "").splitlines() or [""]:
        p = tx.text_frame.add_paragraph() if tx.text_frame.text else tx.text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(16); p.font.color.rgb = DISCIPLINE


def _add_metric_slide(prs, title: str, metrics: list[tuple[str, str]],
                      logo_path: Path | None = None):
    """metrics: list of (label, value) pairs — rendered as big cards."""
    slide = _add_blank_slide(prs)
    _add_header(slide, title, logo_path)
    n = len(metrics)
    if n == 0:
        return
    card_w  = (12.5 - 0.4 * (n - 1)) / n
    for i, (label, value) in enumerate(metrics):
        left = Inches(0.4 + i * (card_w + 0.4))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      left, Inches(2.3), Inches(card_w), Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb = STAMINA
        box.line.color.rgb = ELITE; box.line.width = Pt(1)

        lab = slide.shapes.add_textbox(left, Inches(2.5), Inches(card_w), Inches(0.5))
        p = lab.text_frame.paragraphs[0]
        p.text = str(label)
        p.font.size = Pt(14); p.font.color.rgb = DISCIPLINE
        p.alignment = 2  # centre

        val = slide.shapes.add_textbox(left, Inches(3.1), Inches(card_w), Inches(1.4))
        pv = val.text_frame.paragraphs[0]
        pv.text = str(value)
        pv.font.size = Pt(44); pv.font.bold = True; pv.font.color.rgb = ELITE
        pv.alignment = 2


def build_pptx(deck_title: str, sections: list[dict],
               subtitle: str = "", logo_path: Path | None = None) -> bytes:
    """Build a 16:9 widescreen deck.

    Each section becomes its own slide (spaced out, not cramped).
    Section dict shape:
      {"title": str, "kind": "table"|"chart"|"text"|"metric", ...kind-specific...}
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_cover(prs, deck_title, subtitle, logo_path)

    for sec in sections:
        kind  = sec.get("kind")
        title = sec.get("title", "")
        if kind == "table":
            _add_table_slide(prs, title, sec.get("df"), sec.get("max_rows", 18), logo_path)
        elif kind == "chart":
            _add_chart_slide(prs, title, sec.get("fig"), logo_path)
        elif kind == "text":
            _add_text_slide(prs, title, sec.get("body", ""), logo_path)
        elif kind == "metric":
            _add_metric_slide(prs, title, sec.get("metrics", []), logo_path)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
